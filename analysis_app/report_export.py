"""
Final Report PPT/PDF export — 한/영 양 언어 지원.

build_report() 결과(dict)를 받아 PPT 또는 PDF bytes로 변환.

사용:
    from analysis_app.report_export import export_pptx, export_pdf
    pptx_bytes = export_pptx(report, lang="ko")  # 또는 "en"
    pdf_bytes  = export_pdf(report,  lang="ko")
"""
from __future__ import annotations

import io
from datetime import date
from typing import Any


# ══════════════════════════════════════════════════════════════════════════════
# i18n — 한국어/영어 라벨 매핑
# ══════════════════════════════════════════════════════════════════════════════

I18N: dict[str, dict[str, str]] = {
    "ko": {
        # ── 제목 ──
        "report_title":      "Final Report — POS Signal Analysis",
        "report_subtitle":   "Alternative Data Intelligence Platform",
        "generated":         "생성일",
        "period":            "분석 기간",
        "companies":         "분석 기업",
        "modules":           "실행 모듈",
        # ── 섹션 ──
        "sec_overview":      "1. 개요",
        "sec_highlights":    "2. 데이터 실측 인사이트",
        "sec_selling":       "3. 데이터 차별점 (Bloomberg/공시 대비)",
        "sec_what_happened": "4. 분석 결과 요약 (What Happened)",
        "sec_what_means":    "5. 의미 해석 (What It Means)",
        "sec_what_to_do":    "6. 다음 액션 (What To Do)",
        "sec_use_case":      "7. 글로벌 기관투자자 활용 시나리오",
        "sec_confidence":    "8. 신뢰도 평가",
        # ── KPI 카드 ──
        "alpha_score":       "ALPHA SCORE",
        "strong_signal":     "강한 신호",
        "neutral_signal":    "중립 신호",
        "weak_signal":       "약한 신호",
        "growth":            "Growth",
        "demand":            "Demand",
        "safety":            "Safety",
        "bonus":             "Bonus",
        # ── 기타 ──
        "data_caption":      "본 섹션은 사용자가 업로드한 실제 데이터의 통계입니다.",
        "selling_caption":   "Bloomberg · FactSet · Refinitiv · 공시(DART) 대비 본 데이터의 정량적 우위 요소.",
        "use_case_caption":  "본 데이터를 운용 워크플로우에 통합할 수 있는 글로벌 투자기관 시나리오.",
        "confidence":        "종합 신뢰도",
        "caveats":           "주의 사항",
        "kpi":               "KPI",
        "vs":                "vs",
        "workflow":          "운용 흐름",
        "case":              "케이스",
        "footer":            "본 리포트는 POS 데이터 기반 소비 신호 분석으로 공식 재무 데이터를 대체하지 않습니다.",
    },
    "en": {
        "report_title":      "Final Report — POS Signal Analysis",
        "report_subtitle":   "Alternative Data Intelligence Platform",
        "generated":         "Generated",
        "period":            "Period",
        "companies":         "Companies",
        "modules":           "Modules Run",
        "sec_overview":      "1. Overview",
        "sec_highlights":    "2. Data Highlights — Actual Statistics",
        "sec_selling":       "3. Selling Points — vs Bloomberg/FactSet/DART",
        "sec_what_happened": "4. What Happened — Analysis Summary",
        "sec_what_means":    "5. What It Means — Interpretation",
        "sec_what_to_do":    "6. What To Do — Next Actions",
        "sec_use_case":      "7. Global Institutional Use Cases",
        "sec_confidence":    "8. Confidence Assessment",
        "alpha_score":       "ALPHA SCORE",
        "strong_signal":     "Strong Signal",
        "neutral_signal":    "Neutral Signal",
        "weak_signal":       "Weak Signal",
        "growth":            "Growth",
        "demand":            "Demand",
        "safety":            "Safety",
        "bonus":             "Bonus",
        "data_caption":      "This section shows statistics from the actual uploaded data.",
        "selling_caption":   "Quantitative advantages of this data vs Bloomberg / FactSet / Refinitiv / DART filings.",
        "use_case_caption":  "Scenarios for integrating this data into global institutional investor workflows.",
        "confidence":        "Overall Confidence",
        "caveats":           "Caveats",
        "kpi":               "KPI",
        "vs":                "vs",
        "workflow":          "Workflow",
        "case":              "Case",
        "footer":            "This report is based on POS consumption data analysis and does not replace official financial data.",
    },
}


def _t(lang: str, key: str) -> str:
    return I18N.get(lang, I18N["ko"]).get(key, key)


# ══════════════════════════════════════════════════════════════════════════════
# PPT Export (python-pptx)
# ══════════════════════════════════════════════════════════════════════════════

def export_pptx(report: dict, lang: str = "ko") -> bytes:
    """build_report() 결과 → PowerPoint 바이트 생성.

    슬라이드 구성:
      1. 표지 — 제목 · 기간 · KPI 요약
      2. Data Highlights (실측)
      3. Selling Points (차별점)
      4. What Happened / What It Means / What To Do
      5. Use Cases (상위 3개)
      6. Confidence + Caveats
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError(
            "python-pptx가 설치되지 않았습니다. "
            "터미널에서 `pip install python-pptx`로 설치 후 다시 실행하세요."
        )

    facts = report.get("facts", {}) or {}
    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # 빈 슬라이드

    # ── 디자인 팔레트 (Midnight Executive) ────────────────────────────────────
    C_NAVY    = RGBColor(0x1E, 0x27, 0x61)
    C_BLUE    = RGBColor(0x1D, 0x4E, 0xD8)
    C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
    C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
    C_RED     = RGBColor(0xDC, 0x26, 0x26)
    C_GRAY    = RGBColor(0x64, 0x74, 0x8B)
    C_DARK    = RGBColor(0x0F, 0x17, 0x2A)
    C_LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

    def _add_textbox(slide, left, top, width, height, text,
                     size=14, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf  = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        if isinstance(text, str):
            text = [text]
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = str(line)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Malgun Gothic" if lang == "ko" else "Calibri"
        return box

    def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line_color
        shape.shadow.inherit = False
        return shape

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — 표지
    # ════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    # 전체 navy 배경
    _add_rect(s1, 0, 0, prs.slide_width, prs.slide_height, C_NAVY)
    # 좌측 accent bar
    _add_rect(s1, Inches(0.6), Inches(1.5), Inches(0.1), Inches(2.0), C_WHITE)
    _add_textbox(s1, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
                 _t(lang, "report_subtitle").upper(),
                 size=12, color=RGBColor(0xCA, 0xDC, 0xFC))
    _add_textbox(s1, Inches(0.9), Inches(1.4), Inches(11), Inches(1.4),
                 _t(lang, "report_title"),
                 size=44, bold=True, color=C_WHITE)
    # 메타
    n_co     = len(facts.get("companies", []) or [])
    n_mod    = len(facts.get("modules_run", []) or [])
    date_s   = facts.get("date_start", "—")
    date_e   = facts.get("date_end", "—")
    meta_line = (
        f"{_t(lang,'period')}: {date_s} ~ {date_e}  ·  "
        f"{_t(lang,'companies')}: {n_co}  ·  "
        f"{_t(lang,'modules')}: {n_mod}"
    )
    _add_textbox(s1, Inches(0.9), Inches(3.3), Inches(11), Inches(0.5),
                 meta_line, size=14, color=RGBColor(0xCA, 0xDC, 0xFC))

    # KPI box (alpha score 강조)
    alpha = facts.get("alpha_score") or 0
    color = C_GREEN if alpha >= 75 else C_AMBER if alpha >= 55 else C_RED
    _add_rect(s1, Inches(0.9), Inches(4.5), Inches(3.5), Inches(2.0), C_WHITE)
    _add_textbox(s1, Inches(0.9), Inches(4.6), Inches(3.5), Inches(0.4),
                 _t(lang, "alpha_score"), size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
    _add_textbox(s1, Inches(0.9), Inches(5.0), Inches(3.5), Inches(1.2),
                 f"{alpha:.0f}", size=72, bold=True, color=color, align=PP_ALIGN.CENTER)
    label_key = "strong_signal" if alpha >= 75 else "neutral_signal" if alpha >= 55 else "weak_signal"
    _add_textbox(s1, Inches(0.9), Inches(6.1), Inches(3.5), Inches(0.4),
                 _t(lang, label_key), size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Sub-scores
    sub_labels = ["growth", "demand", "safety", "bonus"]
    sub_keys   = ["growth_pts", "demand_pts", "safety_pts", "bonus_pts"]
    sub_max    = [40, 35, 25, 10]
    for i, (lk, sk, smax) in enumerate(zip(sub_labels, sub_keys, sub_max)):
        x = Inches(4.8 + i * 1.95)
        _add_rect(s1, x, Inches(4.5), Inches(1.8), Inches(2.0), C_WHITE)
        v = facts.get(sk, 0) or 0
        _add_textbox(s1, x, Inches(4.7), Inches(1.8), Inches(0.4),
                     _t(lang, lk), size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(s1, x, Inches(5.1), Inches(1.8), Inches(0.9),
                     f"{v:.0f}", size=42, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        _add_textbox(s1, x, Inches(6.0), Inches(1.8), Inches(0.4),
                     f"/ {smax}", size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

    _add_textbox(s1, Inches(0.9), Inches(6.9), Inches(11), Inches(0.3),
                 f"{_t(lang,'generated')}: {report.get('generated_at', date.today().isoformat())}",
                 size=10, color=RGBColor(0xCA, 0xDC, 0xFC))

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Data Highlights
    # ════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    _add_textbox(s2, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 _t(lang, "sec_highlights"), size=28, bold=True, color=C_NAVY)
    _add_textbox(s2, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                 _t(lang, "data_caption"), size=12, color=C_GRAY)

    dh_items = (report.get("data_highlights", {}) or {}).get("items", [])
    # 2열 grid
    for i, item in enumerate(dh_items[:8]):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.7 + row * 1.35)
        _add_rect(s2, x, y, Inches(6.0), Inches(1.25), C_LIGHT)
        # 좌측 accent
        _add_rect(s2, x, y, Inches(0.08), Inches(1.25), C_BLUE)
        # 제목
        title_text = f"{item.get('icon','')} {item.get('title','')}"
        _add_textbox(s2, Inches(x.inches + 0.2), y, Inches(5.7), Inches(0.4),
                     title_text, size=12, bold=True, color=C_DARK)
        # 본문 (HTML 태그 제거)
        body = str(item.get("body", "")).replace("<b>", "").replace("</b>", "")
        _add_textbox(s2, Inches(x.inches + 0.2), Inches(y.inches + 0.4),
                     Inches(5.7), Inches(0.85),
                     body, size=10, color=C_DARK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Selling Points
    # ════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    _add_textbox(s3, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 _t(lang, "sec_selling"), size=28, bold=True, color=C_NAVY)
    _add_textbox(s3, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                 _t(lang, "selling_caption"), size=12, color=C_GRAY)

    sp_points = (report.get("selling_points", {}) or {}).get("points", [])
    for i, point in enumerate(sp_points[:6]):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.7 + row * 1.85)
        _add_rect(s3, x, y, Inches(6.0), Inches(1.75), C_LIGHT)
        _add_rect(s3, x, y, Inches(0.08), Inches(1.75), C_AMBER)
        ttl = f"{point.get('icon','')} {point.get('title','')}"
        _add_textbox(s3, Inches(x.inches + 0.2), y, Inches(5.7), Inches(0.4),
                     ttl, size=12, bold=True, color=C_DARK)
        hl = point.get("headline", "")
        _add_textbox(s3, Inches(x.inches + 0.2), Inches(y.inches + 0.4),
                     Inches(5.7), Inches(0.4),
                     hl, size=11, bold=True, color=C_AMBER)
        # KPI 라인
        kpi = f"📊 {point.get('kpi', '')}"
        _add_textbox(s3, Inches(x.inches + 0.2), Inches(y.inches + 0.85),
                     Inches(5.7), Inches(0.4),
                     kpi, size=10, color=C_BLUE)
        vs = f"⚖️ vs {point.get('vs', '')}"
        _add_textbox(s3, Inches(x.inches + 0.2), Inches(y.inches + 1.25),
                     Inches(5.7), Inches(0.4),
                     vs, size=10, color=C_GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Findings (What Happened + What It Means + What To Do)
    # ════════════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank_layout)
    _add_textbox(s4, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 _t(lang, "sec_what_happened") + "  ·  "
                 + _t(lang, "sec_what_means").split(". ", 1)[-1] + "  ·  "
                 + _t(lang, "sec_what_to_do").split(". ", 1)[-1],
                 size=22, bold=True, color=C_NAVY)

    sections = [
        ("what_happened",  "📊", C_BLUE),
        ("what_it_means",  "🔍", C_GREEN),
        ("what_to_do",     "🎯", C_AMBER),
    ]
    for i, (key, icon, color) in enumerate(sections):
        data    = report.get(key, {}) or {}
        bullets = data.get("bullets", []) or []
        x = Inches(0.5 + i * 4.3)
        y = Inches(1.3)
        _add_rect(s4, x, y, Inches(4.1), Inches(5.7), C_LIGHT)
        _add_rect(s4, x, y, Inches(0.08), Inches(5.7), color)
        # 섹션 헤더
        head_key = "sec_" + key
        _add_textbox(s4, Inches(x.inches + 0.2), y, Inches(3.85), Inches(0.5),
                     f"{icon} {_t(lang, head_key).split('. ', 1)[-1]}",
                     size=14, bold=True, color=C_DARK)
        # bullets
        lines = []
        for b in bullets[:6]:
            clean = str(b).replace("**", "").replace("*", "")
            if len(clean) > 110:
                clean = clean[:107] + "..."
            lines.append("• " + clean)
        _add_textbox(s4, Inches(x.inches + 0.2), Inches(y.inches + 0.55),
                     Inches(3.85), Inches(5.05),
                     lines, size=10, color=C_DARK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Use Cases (Top 3)
    # ════════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank_layout)
    _add_textbox(s5, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 _t(lang, "sec_use_case"), size=28, bold=True, color=C_NAVY)
    _add_textbox(s5, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                 _t(lang, "use_case_caption"), size=12, color=C_GRAY)

    uc_sections = (report.get("use_case", {}) or {}).get("sections", [])
    for i, sec in enumerate(uc_sections[:3]):
        x = Inches(0.5 + i * 4.3)
        y = Inches(1.7)
        _add_rect(s5, x, y, Inches(4.1), Inches(5.4), C_LIGHT)
        _add_rect(s5, x, y, Inches(0.08), Inches(5.4), C_NAVY)
        _add_textbox(s5, Inches(x.inches + 0.2), y, Inches(3.85), Inches(0.6),
                     sec.get("audience", ""), size=12, bold=True, color=C_DARK)
        _add_textbox(s5, Inches(x.inches + 0.2), Inches(y.inches + 0.55),
                     Inches(3.85), Inches(0.4),
                     sec.get("tagline", ""), size=9, color=C_GRAY)
        value = str(sec.get("value", "")).replace("**", "")
        if len(value) > 280:
            value = value[:277] + "..."
        _add_textbox(s5, Inches(x.inches + 0.2), Inches(y.inches + 1.05),
                     Inches(3.85), Inches(2.4),
                     value, size=10, color=C_DARK)
        _add_textbox(s5, Inches(x.inches + 0.2), Inches(y.inches + 3.55),
                     Inches(3.85), Inches(0.4),
                     f"📊 {sec.get('kpi', '')}", size=10, bold=True, color=C_BLUE)
        pattern = sec.get("use_pattern", "")
        _add_textbox(s5, Inches(x.inches + 0.2), Inches(y.inches + 4.0),
                     Inches(3.85), Inches(1.3),
                     f"📌 {_t(lang, 'workflow')}: {pattern}",
                     size=9, color=C_GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Confidence + Caveats
    # ════════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank_layout)
    _add_textbox(s6, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 _t(lang, "sec_confidence"), size=28, bold=True, color=C_NAVY)

    conf = report.get("confidence", {}) or {}
    overall = conf.get("overall", "—")
    _add_textbox(s6, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                 f"{_t(lang, 'confidence')}: {overall}",
                 size=18, bold=True, color=C_BLUE)

    # Confidence factors 표
    factors = conf.get("factors", []) or []
    for i, f in enumerate(factors[:6]):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.9 + row * 0.9)
        _add_rect(s6, x, y, Inches(6.0), Inches(0.8), C_LIGHT)
        _add_textbox(s6, Inches(x.inches + 0.15), Inches(y.inches + 0.08),
                     Inches(2.5), Inches(0.3),
                     f.get("name", ""), size=11, color=C_GRAY)
        _add_textbox(s6, Inches(x.inches + 0.15), Inches(y.inches + 0.4),
                     Inches(2.5), Inches(0.4),
                     f"{f.get('status','')} {f.get('value','')}",
                     size=13, bold=True, color=C_DARK)
        _add_textbox(s6, Inches(x.inches + 2.8), Inches(y.inches + 0.2),
                     Inches(3.1), Inches(0.5),
                     str(f.get("note", "")), size=10, color=C_DARK)

    # Caveats
    caveats = conf.get("caveats", []) or []
    y_cav = Inches(5.0)
    _add_textbox(s6, Inches(0.5), y_cav, Inches(12), Inches(0.4),
                 f"⚠️ {_t(lang, 'caveats')}", size=13, bold=True, color=C_AMBER)
    cav_lines = ["• " + c for c in caveats[:5]]
    _add_textbox(s6, Inches(0.5), Inches(y_cav.inches + 0.4),
                 Inches(12), Inches(2.0),
                 cav_lines, size=9, color=C_DARK)

    # Footer
    _add_textbox(s6, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
                 _t(lang, "footer"), size=8, color=C_GRAY)

    # ── 출력 ─────────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# PDF Export (reportlab) — A4 portrait, 한글 폰트 자동 시도
# ══════════════════════════════════════════════════════════════════════════════

def _register_kor_font() -> str:
    """시스템에서 사용 가능한 한글 폰트 등록 시도. 성공 시 폰트명 반환.

    macOS / Linux의 다양한 표준 경로 모두 탐색.
    """
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    except Exception:
        return "Helvetica"

    import os, glob

    # 1) reportlab 내장 CID 한글 폰트 — TTF 파일 없어도 동작 (가장 안정적)
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))
        return "HYSMyeongJo-Medium"
    except Exception:
        pass

    # 2) 시스템에서 TTF/OTF 한글 폰트 탐색
    kor_candidates: list[tuple[str, str]] = []   # (font_name, file_path)

    # macOS standard 경로들 (Big Sur+)
    mac_paths = [
        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        "/Library/Fonts/AppleGothic.ttf",
        "/Library/Fonts/AppleSDGothicNeo.ttc",
        "/System/Library/Fonts/Apple SD Gothic Neo.ttc",
        "/Library/Fonts/NanumGothic.ttf",
        "/Library/Fonts/NanumGothic.otf",
    ]
    # Linux standard 경로들
    linux_paths = [
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc",
    ]
    for p in mac_paths + linux_paths:
        if os.path.exists(p):
            base = os.path.splitext(os.path.basename(p))[0].replace(" ", "")
            kor_candidates.append((base, p))

    # 3) 사용자 홈 + ~/Library/Fonts/ 에서 한글 가능성 있는 폰트 탐색
    home = os.path.expanduser("~")
    glob_patterns = [
        os.path.join(home, "Library/Fonts/*Gothic*.ttf"),
        os.path.join(home, "Library/Fonts/*Gothic*.otf"),
        os.path.join(home, "Library/Fonts/*Nanum*.ttf"),
        os.path.join(home, "Library/Fonts/*Noto*Sans*CJK*.ttc"),
        "/System/Library/Fonts/Supplemental/*Gothic*.ttf",
    ]
    for pat in glob_patterns:
        for p in glob.glob(pat):
            base = os.path.splitext(os.path.basename(p))[0].replace(" ", "")
            kor_candidates.append((base, p))

    # 4) 시도 — ttc는 subfontIndex 0~3까지 시도
    seen: set[str] = set()
    for name, path in kor_candidates:
        if path in seen:
            continue
        seen.add(path)
        for idx in (0, 1, 2, 3):
            try:
                pdfmetrics.registerFont(TTFont(name, path, subfontIndex=idx))
                return name
            except Exception:
                continue

    # 모두 실패 — Helvetica로 fallback (한글 ▫▫ 표시됨)
    return "Helvetica"


def export_pdf(report: dict, lang: str = "ko") -> bytes:
    """build_report() 결과 → PDF 바이트 생성 (A4, reportlab Platypus)."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle,
        )
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units  import mm
        from reportlab.lib        import colors
        from reportlab.lib.enums  import TA_LEFT, TA_CENTER
    except ImportError:
        raise RuntimeError(
            "reportlab이 설치되지 않았습니다. "
            "터미널에서 `pip install reportlab`로 설치 후 다시 실행하세요."
        )

    facts    = report.get("facts", {}) or {}
    # 한글 폰트는 양 언어 모두에 등록 — 영문 리포트에 회사명 등 한국어가 남아도 깨지지 않게
    font     = _register_kor_font()
    font_b   = font  # bold variant — fallback to same

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=15 * mm, rightMargin=15 * mm,
        topMargin=15 * mm, bottomMargin=15 * mm,
    )

    base = getSampleStyleSheet()
    st_title = ParagraphStyle(
        "TitleK", parent=base["Title"], fontName=font, fontSize=22,
        textColor=colors.HexColor("#1E2761"), alignment=TA_LEFT, spaceAfter=4,
    )
    st_subtitle = ParagraphStyle(
        "SubK", parent=base["Normal"], fontName=font, fontSize=10,
        textColor=colors.HexColor("#64748B"), spaceAfter=12,
    )
    st_h1 = ParagraphStyle(
        "H1K", parent=base["Heading1"], fontName=font, fontSize=14,
        textColor=colors.HexColor("#1E2761"), spaceAfter=6, spaceBefore=10,
    )
    st_body = ParagraphStyle(
        "BodyK", parent=base["Normal"], fontName=font, fontSize=10,
        textColor=colors.HexColor("#1F2937"), spaceAfter=4, leading=14,
    )
    st_caption = ParagraphStyle(
        "CapK", parent=base["Normal"], fontName=font, fontSize=8,
        textColor=colors.HexColor("#64748B"), spaceAfter=8,
    )

    story = []

    # ── 표지 정보 ────────────────────────────────────────────────────────────
    story.append(Paragraph(_t(lang, "report_subtitle").upper(), st_caption))
    story.append(Paragraph(_t(lang, "report_title"), st_title))
    meta = (
        f"{_t(lang,'period')}: {facts.get('date_start','—')} ~ {facts.get('date_end','—')}  ·  "
        f"{_t(lang,'companies')}: {len(facts.get('companies', []) or [])}  ·  "
        f"{_t(lang,'modules')}: {len(facts.get('modules_run', []) or [])}  ·  "
        f"{_t(lang,'generated')}: {report.get('generated_at', date.today().isoformat())}"
    )
    story.append(Paragraph(meta, st_subtitle))

    # ── Alpha Score Hero ─────────────────────────────────────────────────────
    alpha = facts.get("alpha_score") or 0
    color = (colors.HexColor("#16A34A") if alpha >= 75
             else colors.HexColor("#D97706") if alpha >= 55
             else colors.HexColor("#DC2626"))
    label_key = "strong_signal" if alpha >= 75 else "neutral_signal" if alpha >= 55 else "weak_signal"
    hero_data = [[
        Paragraph(
            f"<font color='{color.hexval()}'><b>{alpha:.0f}</b></font><br/>"
            f"<font size='8' color='#64748B'>{_t(lang, 'alpha_score')} · "
            f"<b>{_t(lang, label_key)}</b></font>",
            ParagraphStyle("hero", fontName=font, fontSize=28, leading=32,
                           textColor=color, alignment=TA_CENTER),
        ),
        Paragraph(
            f"<font color='#1E40AF'><b>{facts.get('growth_pts',0):.0f}/40</b></font> "
            f"{_t(lang,'growth')}<br/>"
            f"<font color='#1E40AF'><b>{facts.get('demand_pts',0):.0f}/35</b></font> "
            f"{_t(lang,'demand')}<br/>"
            f"<font color='#1E40AF'><b>{facts.get('safety_pts',0):.0f}/25</b></font> "
            f"{_t(lang,'safety')}<br/>"
            f"<font color='#1E40AF'><b>{facts.get('bonus_pts',0):.0f}/10</b></font> "
            f"{_t(lang,'bonus')}",
            st_body,
        ),
    ]]
    hero_tbl = Table(hero_data, colWidths=[55 * mm, 110 * mm])
    hero_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (0,0), colors.HexColor("#F8FAFC")),
        ("BACKGROUND", (1,0), (1,0), colors.HexColor("#F8FAFC")),
        ("VALIGN",     (0,0), (-1,-1), "MIDDLE"),
        ("LEFTPADDING",  (0,0), (-1,-1), 12),
        ("RIGHTPADDING", (0,0), (-1,-1), 12),
        ("TOPPADDING",   (0,0), (-1,-1), 14),
        ("BOTTOMPADDING",(0,0), (-1,-1), 14),
        ("ROUNDEDCORNERS", [4,4,4,4]),
    ]))
    story.append(hero_tbl)
    story.append(Spacer(1, 8))

    # ── Data Highlights ──────────────────────────────────────────────────────
    story.append(Paragraph(_t(lang, "sec_highlights"), st_h1))
    story.append(Paragraph(_t(lang, "data_caption"), st_caption))
    for item in (report.get("data_highlights", {}) or {}).get("items", [])[:10]:
        title = f"{item.get('icon','')} <b>{item.get('title','')}</b>"
        body  = str(item.get("body", "")).replace("<b>", "<b>").replace("</b>", "</b>")
        story.append(Paragraph(f"{title} — {body}", st_body))

    # ── Selling Points ───────────────────────────────────────────────────────
    story.append(PageBreak())
    story.append(Paragraph(_t(lang, "sec_selling"), st_h1))
    story.append(Paragraph(_t(lang, "selling_caption"), st_caption))
    for p in (report.get("selling_points", {}) or {}).get("points", []):
        story.append(Paragraph(
            f"{p.get('icon','')} <b>{p.get('title','')}</b>", st_body,
        ))
        story.append(Paragraph(
            f"<font color='#92400E'><b>{p.get('headline','')}</b></font>",
            ParagraphStyle("hl", parent=st_body, leftIndent=10, spaceAfter=2),
        ))
        for d in p.get("details", [])[:3]:
            story.append(Paragraph(
                f"• {d}", ParagraphStyle("dt", parent=st_body, leftIndent=20,
                                          fontSize=9, spaceAfter=1),
            ))
        story.append(Paragraph(
            f"📊 {p.get('kpi','')}  ·  ⚖️ vs {p.get('vs','')}",
            ParagraphStyle("kpi", parent=st_caption, leftIndent=10, spaceAfter=8),
        ))

    # ── What Happened / Means / To Do ───────────────────────────────────────
    story.append(PageBreak())
    for key in ("what_happened", "what_it_means", "what_to_do"):
        sec_key = "sec_" + key
        story.append(Paragraph(_t(lang, sec_key), st_h1))
        for b in (report.get(key, {}) or {}).get("bullets", [])[:8]:
            clean = str(b).replace("**", "<b>", 1).replace("**", "</b>", 1)
            story.append(Paragraph("• " + clean, st_body))
        story.append(Spacer(1, 6))

    # ── Use Cases ────────────────────────────────────────────────────────────
    story.append(PageBreak())
    story.append(Paragraph(_t(lang, "sec_use_case"), st_h1))
    story.append(Paragraph(_t(lang, "use_case_caption"), st_caption))
    for sec in (report.get("use_case", {}) or {}).get("sections", [])[:8]:
        story.append(Paragraph(
            f"<b>{sec.get('audience','')}</b>  "
            f"<font size='8' color='#64748B'><i>{sec.get('tagline','')}</i></font>",
            st_body,
        ))
        story.append(Paragraph(
            str(sec.get("value", "")),
            ParagraphStyle("uc_val", parent=st_body, leftIndent=10, spaceAfter=2),
        ))
        story.append(Paragraph(
            f"📊 {sec.get('kpi','')}  ·  📌 {_t(lang,'workflow')}: {sec.get('use_pattern','')}",
            ParagraphStyle("uc_meta", parent=st_caption, leftIndent=10, spaceAfter=6),
        ))

    # ── Confidence ───────────────────────────────────────────────────────────
    story.append(PageBreak())
    story.append(Paragraph(_t(lang, "sec_confidence"), st_h1))
    conf = report.get("confidence", {}) or {}
    story.append(Paragraph(
        f"<b>{_t(lang, 'confidence')}: {conf.get('overall', '—')}</b>",
        ParagraphStyle("oc", parent=st_body, fontSize=12, textColor=colors.HexColor("#1E40AF")),
    ))
    for f in conf.get("factors", []) or []:
        story.append(Paragraph(
            f"{f.get('status','')} <b>{f.get('name','')}</b>: {f.get('value','')} — {f.get('note','')}",
            st_body,
        ))
    story.append(Spacer(1, 6))
    story.append(Paragraph(f"⚠️ <b>{_t(lang, 'caveats')}</b>", st_body))
    for c in (conf.get("caveats", []) or [])[:6]:
        story.append(Paragraph(f"• {c}",
                                ParagraphStyle("cav", parent=st_caption, leftIndent=10)))

    story.append(Spacer(1, 20))
    story.append(Paragraph(_t(lang, "footer"), st_caption))

    doc.build(story)
    return buf.getvalue()
